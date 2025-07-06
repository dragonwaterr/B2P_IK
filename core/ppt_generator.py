from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import os
import json
import shutil
from copy import deepcopy
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import platform
import subprocess
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

BIBLE_DATA_PATH = os.path.join("data", "cache", "bible_data.json")

# 성경책 약칭 테이블 (약칭: 정식이름)
BIBLE_BOOK_ABBR = {
    "창": "창세기", "출": "출애굽기", "레": "레위기", "민": "민수기", "신": "신명기", "수": "여호수아", "삿": "사사기",
    "룻": "룻기", "삼상": "사무엘상", "삼하": "사무엘하", "왕상": "열왕기상", "왕하": "열왕기하", "대상": "역대상", "대하": "역대하",
    "스": "에스라", "느": "느헤미야", "에": "에스더", "욥": "욥기", "시": "시편", "잠": "잠언", "전": "전도서", "아": "아가",
    "사": "이사야", "렘": "예레미야", "애": "예레미야애가", "겔": "에스겔", "단": "다니엘", "호": "호세아", "욜": "요엘",
    "암": "아모스", "옵": "오바댜", "욘": "요나", "미": "미가", "나": "나훔", "합": "하박국", "습": "스바냐", "학": "학개",
    "슥": "스가랴", "말": "말라기", "마": "마태복음", "막": "마가복음", "눅": "누가복음", "요": "요한복음",
    "행": "사도행전", "롬": "로마서", "고전": "고린도전서", "고후": "고린도후서", "갈": "갈라디아서", "엡": "에베소서",
    "빌": "빌립보서", "골": "골로새서", "살전": "데살로니가전서", "살후": "데살로니가후서", "딤전": "디모데전서", "딤후": "디모데후서",
    "딛": "디도서", "몬": "빌레몬서", "히": "히브리서", "약": "야고보서", "벧전": "베드로전서", "벧후": "베드로후서",
    "요일": "요한1서", "요이": "요한2서", "요삼": "요한3서", "유": "유다서", "계": "요한계시록"
}

def load_bible_data():
    with open(BIBLE_DATA_PATH, "r", encoding="utf-8") as f:
        return json.load(f)

def parse_selection(selection_str, bible_data, version):
    """
    예시 입력: '창세기1:1-3,2:1-2; 출애굽기3:1-5' 또는 '창1:1-3,2:1-2; 왕상3:1-5'
    또는 '창4:1-3,2:1-2; 왕상3:1-5' (공백 없이도 인식)
    여러 권, 여러 범위 지원. 반환: [(책, 장, [절, ...]), ...]
    """
    result = []
    for book_part in selection_str.split(";"):
        book_part = book_part.strip()
        if not book_part:
            continue
        
        # 성경이름과 나머지 부분을 분리
        # 숫자가 나오는 첫 번째 위치를 찾아서 분리
        match = re.search(r'(\d+:)', book_part)
        if not match:
            continue
        
        # 숫자 앞까지가 성경이름, 숫자부터가 장절 정보
        split_pos = match.start()
        book_name = book_part[:split_pos].strip()
        rest = book_part[split_pos:]
        
        # 약칭이면 fullname으로 변환
        if book_name in BIBLE_BOOK_ABBR:
            book_name_full = BIBLE_BOOK_ABBR[book_name]
        else:
            book_name_full = book_name
            
        for chapter_range in rest.split(","):
            chapter_range = chapter_range.strip()
            if not chapter_range:
                continue
            if ":" not in chapter_range:
                continue
            chapter, verses = chapter_range.split(":")
            chapter = int(chapter)
            for v in verses.split(","):
                v = v.strip()
                if not v:
                    continue
                if "-" in v:
                    start, end = map(int, v.split("-"))
                    verse_list = list(range(start, end+1))
                else:
                    verse_list = [int(v)]
                result.append((book_name_full, chapter, verse_list))
    return result

def get_verses(bible_data, version, selections):
    version = "개역개정"  # 항상 개역개정만 사용
    verses = []
    for book, chapter, verse_list in selections:
        for verse in verse_list:
            try:
                text = bible_data[version][book][str(chapter)][str(verse)]
            except KeyError:
                continue
            verses.append((book, chapter, verse, text))
    return verses

def clone_slide(prs, slide):
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    # 기존 새 슬라이드의 요소를 모두 삭제
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    # 샘플 슬라이드의 모든 요소 복제
    for shape in slide.shapes:
        el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
    return new_slide

def set_text_preserve_style(shape, new_text):
    """텍스트박스의 스타일을 유지하며 텍스트만 변경"""
    tf = shape.text_frame
    # 기존 paragraph/run 구조 유지, 첫 run만 텍스트 교체
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
        # 나머지 run은 모두 삭제
        for run in tf.paragraphs[0].runs[1:]:
            p = tf.paragraphs[0]._element
            p.remove(run._r)
        # 나머지 paragraph도 모두 삭제
        for para in tf.paragraphs[1:]:
            tf._element.remove(para._p)
    else:
        # 스타일이 없는 경우 fallback
        shape.text = new_text

def copy_slide_background(src_slide, dst_slide):
    src_fill = src_slide.background.fill
    dst_fill = dst_slide.background.fill
    if src_fill.type == 'solid':
        dst_fill.solid()
        dst_fill.fore_color.rgb = src_fill.fore_color.rgb
        # 투명도도 복제
        dst_fill.fore_color.transparency = src_fill.fore_color.transparency
    elif src_fill.type == 'gradient':
        # gradient 등은 복잡하므로, 필요시 추가 구현
        print("경고: 템플릿 배경이 그라데이션입니다. 현재는 단색만 복제합니다.")
    elif src_fill.type == 'patterned':
        print("경고: 템플릿 배경이 패턴입니다. 현재는 단색만 복제합니다.")
    # 기타 fill type은 필요시 추가 구현

def create_ppt(verses, output_path, template_path, bg_image_path=None, max_chars=500):
    shutil.copy(template_path, output_path)
    prs = Presentation(output_path)
    template_slide = prs.slides[0]

    if bg_image_path and os.path.exists(bg_image_path):
        # 배경 투명 + 이미지 첨부
        template_slide.background.fill.solid()
        template_slide.background.fill.transparency = 1.0
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        pic = template_slide.shapes.add_picture(bg_image_path, 0, 0, slide_width, slide_height)
        template_slide.shapes._spTree.remove(pic._element)
        template_slide.shapes._spTree.insert(2, pic._element)
        print("배경 이미지가 존재하는 템플릿으로 생성!!!")
    else:
        # 검정색 배경
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        shape = template_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0,0,0)
        shape.line.fill.background()  # 테두리 없음
        # 맨 뒤로 이동
        template_slide.shapes._spTree.remove(shape._element)
        template_slide.shapes._spTree.insert(2, shape._element)
        print("배경 이미지가 없어요 검정배경으로 !!!")

    # 기존 슬라이드(샘플)는 모두 삭제
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    required_names = {"TitleBox", "VerseBox", "ContentBox"}
    actual_names = {getattr(shape, "name", "") for shape in template_slide.shapes if shape.has_text_frame}
    missing = required_names - actual_names
    if missing:
        raise Exception(f"템플릿에 다음 텍스트박스 이름이 없습니다: {missing}")

    # 입력 순서대로 각 구절별로 슬라이드 생성 (중복 제거하지 않음)
    for book, chapter, verse, text in verses:
        # 텍스트가 최대 글자 수를 초과하면 분할
        if len(text) <= max_chars:
            # 한 슬라이드에 모두 표시
            _create_slide(prs, template_slide, book, chapter, verse, text, bg_image_path)
        else:
            # 여러 슬라이드로 분할
            _create_split_slides(prs, template_slide, book, chapter, verse, text, max_chars, bg_image_path)
    
    prs.save(output_path)

def _create_slide(prs, template_slide, book, chapter, verse, text, bg_image_path):
    """단일 슬라이드 생성"""
    slide = clone_slide(prs, template_slide)
    
    # 배경 이미지 설정
    if bg_image_path and os.path.exists(bg_image_path):
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        pic = slide.shapes.add_picture(bg_image_path, 0, 0, slide_width, slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    
    # 텍스트 설정
    name_to_text = {
        "TitleBox": f"{book} {chapter}장",
        "VerseBox": str(verse),
        "ContentBox": text
    }
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        shape_name = getattr(shape, "name", "")
        if shape_name in name_to_text:
            set_text_preserve_style(shape, name_to_text[shape_name])

def _create_split_slides(prs, template_slide, book, chapter, verse, full_text, max_chars, bg_image_path):
    """긴 텍스트를 여러 슬라이드로 분할"""
    words = full_text.split()
    current_text = ""
    slide_count = 0
    
    for word in words:
        test_text = current_text + " " + word if current_text else word
        if len(test_text) <= max_chars:
            current_text = test_text
        else:
            # 현재 슬라이드 저장
            if current_text:
                slide_count += 1
                verse_display = f"{verse}" if slide_count > 1 else str(verse)
                _create_slide(prs, template_slide, book, chapter, verse_display, current_text, bg_image_path)
            current_text = word
    
    # 마지막 슬라이드 저장
    if current_text:
        slide_count += 1
        verse_display = f"{verse}" if slide_count > 1 else str(verse)
        _create_slide(prs, template_slide, book, chapter, verse_display, current_text, bg_image_path)